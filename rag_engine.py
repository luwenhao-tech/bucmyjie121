"""RAG 引擎（轻量版）：PDF/Excel 论文解析 + 关键词检索。

设计目标：在低内存服务器（< 1GB）上也能运行。
- 索引构建：在本地运行，生成 JSON 索引文件
- 检索：加载 JSON，用 jieba 分词 + BM25 评分，内存占用极低

使用方式：
1. 本地：把论文放入 papers/ 目录，运行 python build_index.py
2. 推送 papers_index.json 到 GitHub
3. 服务器 git pull 后重启即可，无需在服务器上构建索引
"""
import os
import json
import math
import re
import hashlib
from pathlib import Path
from typing import List, Dict, Optional
from collections import Counter

# ============ 配置 ============
PAPERS_DIR = os.getenv("PAPERS_DIR", "papers")
INDEX_FILE = os.getenv("RAG_INDEX_FILE", "papers_index.json")
CHUNK_SIZE = int(os.getenv("RAG_CHUNK_SIZE", "500"))
CHUNK_OVERLAP = int(os.getenv("RAG_CHUNK_OVERLAP", "100"))
TOP_K = int(os.getenv("RAG_TOP_K", "5"))


def ensure_dirs():
    os.makedirs(PAPERS_DIR, exist_ok=True)


# ============ PDF 解析 ============
def extract_text_from_pdf(pdf_path: str) -> str:
    """从 PDF 提取全部文本"""
    import fitz
    doc = fitz.open(pdf_path)
    text_parts = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text("text")
        if text.strip():
            text_parts.append(text)
    doc.close()
    return "\n".join(text_parts)


def extract_metadata_from_pdf(pdf_path: str) -> Dict[str, str]:
    import fitz
    doc = fitz.open(pdf_path)
    metadata = doc.metadata or {}
    doc.close()
    return {
        "title": metadata.get("title", "") or Path(pdf_path).stem,
        "author": metadata.get("author", ""),
        "filename": Path(pdf_path).name,
    }


# ============ Word (.docx) 解析 ============
def extract_text_from_docx(docx_path: str) -> str:
    """从 .docx 提取全部文本（段落 + 表格单元格）"""
    from docx import Document
    doc = Document(docx_path)
    parts = []
    # 正文段落
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # 表格内容
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    parts.append(cell_text)
    return "\n".join(parts)


# ============ PowerPoint (.pptx) 解析 ============
def extract_text_from_pptx(pptx_path: str) -> str:
    """从 .pptx 提取全部文本（每张幻灯片的文本框 + 表格单元格 + 备注）"""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    parts = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_parts = [f"[幻灯片 {slide_idx}]"]
        for shape in slide.shapes:
            # 文本框
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        slide_parts.append(txt)
            # 表格
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            slide_parts.append(cell_text)
        # 演讲者备注
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip() if slide.notes_slide.notes_text_frame else ""
            if notes:
                slide_parts.append(f"[备注] {notes}")
        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))
    return "\n\n".join(parts)


# ============ Excel 解析 ============
def extract_papers_from_excel(xlsx_path: str) -> List[Dict[str, str]]:
    import openpyxl
    wb = openpyxl.load_workbook(xlsx_path, read_only=True)
    papers = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if len(rows) < 2:
            continue
        header = [str(cell or "").strip() for cell in rows[0]]
        col_map = {}
        for i, h in enumerate(header):
            h_lower = h.lower()
            if "论文" in h or "名称" in h or "标题" in h or "title" in h_lower:
                col_map["title"] = i
            elif "作者" in h or "author" in h_lower:
                col_map["author"] = i
            elif "方法" in h or "method" in h_lower:
                col_map["method"] = i
            elif "关键" in h or "keyword" in h_lower:
                col_map["keywords"] = i
            elif "摘要" in h or "abstract" in h_lower:
                col_map["abstract"] = i
            elif "期刊" in h or "journal" in h_lower:
                col_map["journal"] = i
            elif "年" in h or "year" in h_lower:
                col_map["year"] = i
        if "title" not in col_map:
            continue
        for row in rows[1:]:
            if not row or not row[col_map.get("title", 0)]:
                continue
            paper = {
                "title": str(row[col_map["title"]] or "").strip(),
                "author": str(row[col_map.get("author", -1)] or "").strip() if "author" in col_map and col_map["author"] < len(row) else "",
                "method": str(row[col_map.get("method", -1)] or "").strip() if "method" in col_map and col_map["method"] < len(row) else "",
                "keywords": str(row[col_map.get("keywords", -1)] or "").strip() if "keywords" in col_map and col_map["keywords"] < len(row) else "",
                "abstract": str(row[col_map.get("abstract", -1)] or "").strip() if "abstract" in col_map and col_map["abstract"] < len(row) else "",
                "journal": str(row[col_map.get("journal", -1)] or "").strip() if "journal" in col_map and col_map["journal"] < len(row) else "",
                "year": str(row[col_map.get("year", -1)] or "").strip() if "year" in col_map and col_map["year"] < len(row) else "",
            }
            if paper["title"]:
                papers.append(paper)
    wb.close()
    return papers


def paper_to_text(paper: Dict[str, str]) -> str:
    parts = []
    parts.append(f"论文标题：{paper['title']}")
    if paper.get("author"):
        parts.append(f"作者：{paper['author']}")
    if paper.get("journal") or paper.get("year"):
        journal_info = paper.get("journal", "")
        if paper.get("year"):
            journal_info += f"（{paper['year']}）"
        parts.append(f"发表于：{journal_info}")
    if paper.get("keywords"):
        parts.append(f"关键词：{paper['keywords']}")
    if paper.get("method"):
        parts.append(f"研究方法：{paper['method']}")
    if paper.get("abstract"):
        parts.append(f"研究内容：{paper['abstract']}")
    return "\n".join(parts)


# ============ 文本切块 ============
def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    if not text or not text.strip():
        return []
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    chunks = []
    current_chunk = ""
    for para in paragraphs:
        if len(para) > chunk_size:
            if current_chunk:
                chunks.append(current_chunk)
                current_chunk = ""
            sentences = _split_sentences(para)
            for sent in sentences:
                if len(current_chunk) + len(sent) <= chunk_size:
                    current_chunk += sent
                else:
                    if current_chunk:
                        chunks.append(current_chunk)
                    current_chunk = sent
        elif len(current_chunk) + len(para) + 1 <= chunk_size:
            current_chunk += ("\n" if current_chunk else "") + para
        else:
            if current_chunk:
                chunks.append(current_chunk)
            current_chunk = para
    if current_chunk:
        chunks.append(current_chunk)
    if overlap > 0 and len(chunks) > 1:
        overlapped_chunks = [chunks[0]]
        for i in range(1, len(chunks)):
            prev_tail = chunks[i - 1][-overlap:]
            overlapped_chunks.append(prev_tail + "\n" + chunks[i])
        chunks = overlapped_chunks
    chunks = [c for c in chunks if len(c.strip()) > 30]
    return chunks


def _split_sentences(text: str) -> List[str]:
    parts = re.split(r'(?<=[。；;！!？?])|(?<=\. )', text)
    return [p for p in parts if p.strip()]


# ============ 轻量分词（不依赖 jieba，用简单正则）============
def tokenize(text: str) -> List[str]:
    """简单中英文分词：中文按字/常用双字，英文按单词"""
    text = text.lower()
    # 提取英文单词
    en_words = re.findall(r'[a-z]{2,}', text)
    # 提取中文：按 2-4 字的滑动窗口生成 n-gram
    cn_chars = re.findall(r'[\u4e00-\u9fff]+', text)
    cn_tokens = []
    for seg in cn_chars:
        # 生成 bigram 和 trigram
        for i in range(len(seg)):
            if i + 2 <= len(seg):
                cn_tokens.append(seg[i:i+2])
            if i + 3 <= len(seg):
                cn_tokens.append(seg[i:i+3])
        # 也加入单字（但权重低，靠频率自然调节）
        for ch in seg:
            cn_tokens.append(ch)
    # 数字/年份：文件名和正文里出现的 "2025""2020""2351" 这类年份/编号也保留
    num_tokens = re.findall(r'\d{2,}', text)
    return en_words + cn_tokens + num_tokens


# ============ chunk 打分用的 token 空间 ============
# 只用 chunk.text 打分会漏掉「关键词只在文件名/标题里」的场景：
# 例如《2025版中国药典—桑白皮.docx》正文全是【性状】【鉴别】条文——
# 「2025」「药典」这些词只在文件名里出现——用户问「2025 版药典桑白皮」
# 时该 chunk 完全命不上，会被其他老论文挤掉。
#
# 修法：把 filename + title 一并 tokenize，与正文 token 合并后一起参与
# BM25 打分。filename token 复制若干次以软加权，让文件名里的关键词
# 有明显竞争力，但不至于完全主导排序、把内容不匹配的旧文件顶上来。
# build_index 的 doc_freq 和 search 的每 chunk tf 都走同一个 helper 保证一致。
FILENAME_META_BOOST = 3   # 文件名 / 标题 token 复制的次数（软加权系数）


def _chunk_tokens(chunk: Dict) -> List[str]:
    """chunk 打分用的 token 列表：正文 + filename + title（后两者软加权）。"""
    tokens = tokenize(chunk.get("text", ""))
    meta_parts = []
    filename = chunk.get("filename", "")
    if filename:
        # 去掉扩展名，避免 "pdf""docx" 等通用后缀污染 doc_freq
        stem = re.sub(r"\.(pdf|PDF|docx|DOCX|xlsx|XLSX|pptx|PPTX|txt)$", "", filename)
        meta_parts.append(stem)
    title = chunk.get("title", "")
    if title and title not in meta_parts:
        meta_parts.append(title)
    if meta_parts:
        meta_tokens = tokenize(" ".join(meta_parts))
        for _ in range(FILENAME_META_BOOST):
            tokens.extend(meta_tokens)
    return tokens


# ============ 英文 PDF → 中文摘要（供 BM25 中文检索）============
# 背景：tokenize() 对英文只切 ASCII 单词，对中文按 n-gram。
# 全英文 PDF 入库后没有任何中文 token，中文 query 完全检索不到。
# 这里给每篇英文 PDF 生成一段中文摘要写到 sidecar 文件，作为额外 chunk 注入索引，
# chunk 的 filename 仍归属原 PDF，所以引用/citation 是正确的。
ZH_SUMMARY_SUFFIX = "_zh.txt"
ZH_SUMMARY_INPUT_CHARS = 2000  # 喂给 LLM 的英文前缀长度（标题+abstract 通常够了）

ZH_SUMMARY_PROMPT = """你是中医药文献的中英翻译/摘要助手。下面是一篇英文论文的开头部分，请用中文产出一个结构化摘要，专门用于中文关键词检索（BM25），所以中文术语要尽量铺开、同义词都列上。

严格按以下三段输出，每段独占一行，不要加多余解释：
标题：<论文的中文译名，化合物/药材英文名一并给出中文常用名，例如 Morus alba → 桑（桑白皮）>
关键词：<8-15 个中文关键词，逗号分隔，覆盖：中药材中文名、化合物中文名、药理活性、研究方法、疾病/靶点>
摘要：<2-3 句中文，说清研究了什么药材的什么部位、用什么方法、得到了什么主要结论，控制在 200 字以内>

输出总长度控制在 300-500 字。只输出三段内容，不要前后客套、不要 Markdown、不要英文原文。"""


def is_english_text(text: str, cjk_ratio_cutoff: float = 0.05) -> bool:
    """CJK 字符占总非空白字符比 < 5% 视为英文。空文本返回 False。"""
    if not text or not text.strip():
        return False
    cjk = sum(1 for ch in text if '一' <= ch <= '鿿')
    total = sum(1 for ch in text if not ch.isspace())
    if total == 0:
        return False
    return (cjk / total) < cjk_ratio_cutoff


def _zh_summary_path(pdf_path: Path) -> Path:
    """papers/foo.pdf → papers/foo_zh.txt"""
    return pdf_path.with_name(pdf_path.stem + ZH_SUMMARY_SUFFIX)


async def _call_llm_for_zh_summary(text: str) -> str:
    """仿 llm_client.classify_intent 的 bare-call 模式，不注入 Liu Chunsheng persona。"""
    # 局部 import 避免循环依赖在模块加载时触发
    from llm_client import client as _llm_client, MODEL as _LLM_MODEL
    resp = await _llm_client.chat.completions.create(
        model=_LLM_MODEL,
        messages=[
            {"role": "system", "content": ZH_SUMMARY_PROMPT},
            {"role": "user", "content": text[:ZH_SUMMARY_INPUT_CHARS]},
        ],
        temperature=0.2,
        max_tokens=600,
    )
    return (resp.choices[0].message.content or "").strip()


def get_or_generate_zh_summary(pdf_path: Path, text: str, force: bool = False) -> Optional[str]:
    """命中缓存 → 读 sidecar；未命中 → 调 LLM → 写 sidecar。失败返回 None。

    注意：内部用 asyncio.run() 调异步 LLM client。仅限 build_index.py CLI 调用，
    不要在已经在运行 event loop 的异步上下文里调用（会抛 RuntimeError）。
    main.py 不 import build_index，已确认安全。
    """
    import asyncio
    sidecar = _zh_summary_path(pdf_path)

    # 缓存命中
    if sidecar.exists() and not force:
        try:
            cached = sidecar.read_text(encoding="utf-8").strip()
            return cached or None
        except Exception:
            pass  # 读不出来就当没缓存，重生成

    # force 时先删掉旧缓存
    if force and sidecar.exists():
        try:
            sidecar.unlink()
        except Exception:
            pass

    if not text or not text.strip():
        return None

    try:
        summary = asyncio.run(_call_llm_for_zh_summary(text))
        if not summary:
            return None
        sidecar.write_text(summary, encoding="utf-8")
        print(f"  [中文摘要] 生成 → {sidecar.name}")
        return summary
    except Exception as e:
        print(f"  [警告] {pdf_path.name} 中文摘要生成失败：{e}（按英文索引）")
        return None


# ============ 可信度评分（按 filename 索引）============
# credibility_scores.json 由 scoring/score_prompt.py 生成
# 结构：{filename: {"credibility": 0.77, "tier": "B", "doc_kind": "...", ...}}
CREDIBILITY_FILE = "scoring/credibility_scores.json"
# A/B/C 分档阈值（跟 scoring/schema.md 的 tier_of 保持一致，A 门槛已下调到 0.80）
CRED_A_THRESHOLD = 0.80
CRED_B_THRESHOLD = 0.60
# BM25 分数与可信度融合：final = bm25 * (BASE + (1-BASE) * credibility)
# BASE=0.4 意味着 credibility=0 的文献仍能拿到原分的 40%（不至于完全消失），
# credibility=1 的文献拿满 100%。低质文献被明显降权但不会被完全过滤。
CRED_WEIGHT_BASE = 0.4

_credibility_data: Optional[Dict[str, Dict]] = None


def load_credibility() -> Dict[str, Dict]:
    """加载 filename → 打分记录 的映射。文件不存在时返回空 dict（等价于全部无加权）。"""
    global _credibility_data
    if _credibility_data is not None:
        return _credibility_data
    path = Path(__file__).parent / CREDIBILITY_FILE
    if not path.exists():
        _credibility_data = {}
        return _credibility_data
    try:
        raw = json.loads(path.read_text(encoding="utf-8"))
        # 只保留有 credibility 字段的（跳过打分失败的 error 记录）
        _credibility_data = {
            fn: rec for fn, rec in raw.items()
            if isinstance(rec, dict) and "credibility" in rec
        }
        print(f"[可信度] 加载 {len(_credibility_data)} 条评分（{CREDIBILITY_FILE}）")
    except Exception as e:
        print(f"[可信度] 加载失败：{e}，按无加权模式运行")
        _credibility_data = {}
    return _credibility_data


def _cred_of(filename: str) -> float:
    """未打分文件返回 0.6（B 类中位数，中性默认，避免新增文件被完全屏蔽）。"""
    rec = load_credibility().get(filename)
    if rec is None:
        return 0.6
    return float(rec.get("credibility", 0.6))


def _tier_of_score(cred: float) -> str:
    """按当前阈值分档，跟 search 时的加权和分层配额保持一致。"""
    if cred >= CRED_A_THRESHOLD:
        return "A"
    if cred >= CRED_B_THRESHOLD:
        return "B"
    return "C"


def _tier_of_filename(filename: str) -> str:
    return _tier_of_score(_cred_of(filename))


# ============ 索引数据（内存中）============
_index_data: Optional[Dict] = None


def _get_index_path() -> str:
    """获取索引文件的完整路径"""
    # 先检查项目根目录
    root_path = Path(__file__).parent / INDEX_FILE
    if root_path.exists():
        return str(root_path)
    # 再检查当前工作目录
    cwd_path = Path(INDEX_FILE)
    if cwd_path.exists():
        return str(cwd_path)
    return str(root_path)


def load_index() -> Dict:
    """加载索引文件到内存"""
    global _index_data
    if _index_data is not None:
        return _index_data

    index_path = _get_index_path()
    if not Path(index_path).exists():
        _index_data = {"chunks": [], "doc_freq": {}, "total_docs": 0}
        return _index_data

    with open(index_path, "r", encoding="utf-8") as f:
        _index_data = json.load(f)
    return _index_data


def save_index(data: Dict):
    """保存索引到 JSON 文件"""
    global _index_data
    index_path = str(Path(__file__).parent / INDEX_FILE)
    with open(index_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False)
    _index_data = data
    print(f"  索引已保存到: {index_path}")


# ============ 索引构建 ============
def build_index(papers_dir: str = PAPERS_DIR, force: bool = False) -> Dict[str, int]:
    """扫描目录下所有 PDF 和 Excel 并构建索引，保存为 JSON。"""
    ensure_dirs()
    papers_path = Path(papers_dir)

    if not papers_path.exists():
        print(f"[错误] 论文目录不存在: {papers_dir}")
        return {}

    pdf_files = list(papers_path.glob("*.pdf")) + list(papers_path.glob("*.PDF"))
    xlsx_files = list(papers_path.glob("*.xlsx")) + list(papers_path.glob("*.XLSX"))
    docx_files = list(papers_path.glob("*.docx")) + list(papers_path.glob("*.DOCX"))
    pptx_files = list(papers_path.glob("*.pptx")) + list(papers_path.glob("*.PPTX"))
    # 只收 _ocr.txt；显式排除 _zh.txt（中文摘要 sidecar，已经在 PDF 循环里以原 PDF 名义注入索引）
    txt_files = [p for p in papers_path.glob("*_ocr.txt") if not p.name.endswith(ZH_SUMMARY_SUFFIX)]

    if not pdf_files and not xlsx_files and not docx_files and not pptx_files and not txt_files:
        print(f"[提示] {papers_dir}/ 目录下没有可处理的文件")
        return {}

    all_chunks = []  # [{"text": ..., "title": ..., "filename": ...}, ...]
    results = {}

    # 处理 PDF
    if pdf_files:
        print(f"\n{'='*50}")
        print(f"处理 PDF：{len(pdf_files)} 个文件")
        print(f"{'='*50}\n")
        for pdf_file in sorted(pdf_files):
            try:
                text = extract_text_from_pdf(str(pdf_file))
                if not text.strip():
                    print(f"  [警告] {pdf_file.name} 文本为空，跳过")
                    results[pdf_file.name] = 0
                    continue
                metadata = extract_metadata_from_pdf(str(pdf_file))

                # 英文 PDF：先尝试生成/读取中文摘要，作为第一个 chunk 注入
                # （filename 仍归属原 PDF，引用正确）
                zh_chunk = None
                if is_english_text(text):
                    zh_summary = get_or_generate_zh_summary(pdf_file, text, force=force)
                    if zh_summary:
                        zh_chunk = {
                            "text": f"【中文摘要】{metadata['title']}\n{zh_summary}",
                            "title": metadata["title"],
                            "filename": metadata["filename"],
                        }

                chunks = chunk_text(text)
                if zh_chunk:
                    all_chunks.append(zh_chunk)
                for chunk in chunks:
                    all_chunks.append({
                        "text": chunk,
                        "title": metadata["title"],
                        "filename": metadata["filename"],
                    })
                total_blocks = len(chunks) + (1 if zh_chunk else 0)
                extra = " (含中文摘要)" if zh_chunk else ""
                print(f"  [完成] {pdf_file.name} → {total_blocks} 个文本块{extra}")
                results[pdf_file.name] = total_blocks
            except Exception as e:
                print(f"  [错误] {pdf_file.name}: {e}")
                results[pdf_file.name] = -1

    # 处理 OCR 文本
    if txt_files:
        print(f"\n{'='*50}")
        print(f"处理 OCR 文本：{len(txt_files)} 个文件")
        print(f"{'='*50}\n")
        for txt_file in sorted(txt_files):
            try:
                with open(txt_file, "r", encoding="utf-8") as f:
                    text = f.read()
                if not text.strip():
                    continue
                chunks = chunk_text(text)
                title = txt_file.stem.replace("_ocr", "")
                # 若同目录存在同名 PDF（例如扫描版 PDF 提取失败、改用 OCR 文本），
                # 把 filename 改写成 PDF 名，让引用链路统一指向 PDF（跟 _zh.txt 的做法一致）
                companion_pdf = txt_file.with_name(title + ".pdf")
                display_filename = companion_pdf.name if companion_pdf.exists() else txt_file.name
                for chunk in chunks:
                    all_chunks.append({
                        "text": chunk,
                        "title": title,
                        "filename": display_filename,
                    })
                extra = " (归原 PDF)" if display_filename != txt_file.name else ""
                print(f"  [完成] {txt_file.name} → {len(chunks)} 个文本块{extra}")
                results[txt_file.name] = len(chunks)
            except Exception as e:
                print(f"  [错误] {txt_file.name}: {e}")
                results[txt_file.name] = -1

    # 处理 Word (.docx)
    if docx_files:
        print(f"\n{'='*50}")
        print(f"处理 Word：{len(docx_files)} 个文件")
        print(f"{'='*50}\n")
        for docx_file in sorted(docx_files):
            try:
                text = extract_text_from_docx(str(docx_file))
                if not text.strip():
                    print(f"  [警告] {docx_file.name} 文本为空，跳过")
                    results[docx_file.name] = 0
                    continue
                title = docx_file.stem
                chunks = chunk_text(text)
                for chunk in chunks:
                    all_chunks.append({
                        "text": chunk,
                        "title": title,
                        "filename": docx_file.name,
                    })
                print(f"  [完成] {docx_file.name} → {len(chunks)} 个文本块")
                results[docx_file.name] = len(chunks)
            except Exception as e:
                print(f"  [错误] {docx_file.name}: {e}")
                results[docx_file.name] = -1

    # 处理 PowerPoint (.pptx)
    if pptx_files:
        print(f"\n{'='*50}")
        print(f"处理 PowerPoint：{len(pptx_files)} 个文件")
        print(f"{'='*50}\n")
        for pptx_file in sorted(pptx_files):
            try:
                text = extract_text_from_pptx(str(pptx_file))
                if not text.strip():
                    print(f"  [警告] {pptx_file.name} 文本为空，跳过")
                    results[pptx_file.name] = 0
                    continue
                title = pptx_file.stem
                chunks = chunk_text(text)
                for chunk in chunks:
                    all_chunks.append({
                        "text": chunk,
                        "title": title,
                        "filename": pptx_file.name,
                    })
                print(f"  [完成] {pptx_file.name} → {len(chunks)} 个文本块")
                results[pptx_file.name] = len(chunks)
            except Exception as e:
                print(f"  [错误] {pptx_file.name}: {e}")
                results[pptx_file.name] = -1

    # 处理 Excel
    if xlsx_files:
        print(f"\n{'='*50}")
        print(f"处理 Excel：{len(xlsx_files)} 个文件")
        print(f"{'='*50}\n")
        for xlsx_file in sorted(xlsx_files):
            try:
                papers = extract_papers_from_excel(str(xlsx_file))
                for paper in papers:
                    text = paper_to_text(paper)
                    if len(text.strip()) > 30:
                        all_chunks.append({
                            "text": text,
                            "title": paper.get("title", ""),
                            "filename": xlsx_file.name,
                        })
                print(f"  [完成] {xlsx_file.name} → {len(papers)} 篇论文摘要")
                results[xlsx_file.name] = len(papers)
            except Exception as e:
                print(f"  [错误] {xlsx_file.name}: {e}")
                results[xlsx_file.name] = -1

    # 计算文档频率（IDF 用）—— 统一用 _chunk_tokens，让 filename/title 也计入
    doc_freq: Dict[str, int] = Counter()
    for chunk in all_chunks:
        tokens = set(_chunk_tokens(chunk))
        for token in tokens:
            doc_freq[token] += 1

    # 保存索引
    index_data = {
        "chunks": all_chunks,
        "doc_freq": dict(doc_freq),
        "total_docs": len(all_chunks),
    }
    save_index(index_data)

    print(f"\n{'='*50}")
    print(f"索引完成：{len(all_chunks)} 个文本块，来自 {len(results)} 个文件")
    print(f"{'='*50}\n")

    return results


# ============ 中文 query → 英文术语扩展（解决英文 PDF 全文检索盲区）============
# 背景：英文 PDF 的全文 chunk 在索引里只有英文 token；
# 用户用纯中文 query 时 BM25 只命中 LLM 生成的中文摘要 chunk（一篇仅 1 个）。
# 为了让中文 query 也能命中英文原文（Methods/Results/数据），
# 在 search 之前把中文 query 翻成同义的英文术语集合，与原 query 拼接后再分词。
QUERY_EXPAND_PROMPT = """你是中医药/天然药物化学文献的中英术语对照助手。
用户的中文检索词如下，请输出 5-10 个**最有可能出现在英文论文里的对应术语**，逗号分隔，仅一行。
要求：
1. 中药材给拉丁学名 + 常用英文（如：桑白皮 → Morus alba, Mori Cortex Radicis, mulberry root bark）
2. 化合物给英文常用名（如：桑酮 → Morusin；黄酮 → flavonoid, flavone）
3. 药理/方法给英文术语（如：凋亡 → apoptosis；抗菌 → antibacterial, antimicrobial）
4. 不要解释、不要中文、不要标点末尾，仅输出术语列表本身
5. 用户词若已经是英文，原样补几个同义词即可"""

# 进程内缓存：中文 query → 英文术语字符串
# key 用 hash 控制内存（极少超 1MB）；进程重启即丢失，可接受
_query_expand_cache: Dict[str, str] = {}
_QUERY_EXPAND_MAX_LEN = 80   # 超长 query 不扩展（句子级，已包含上下文）
_QUERY_EXPAND_TIMEOUT = 4.0  # LLM 翻译超时秒数，超时则原样检索


def _is_mostly_chinese(text: str, cjk_ratio_cutoff: float = 0.30) -> bool:
    """CJK 占比 >= 30% 视为"以中文为主"，需要做英文扩展。"""
    if not text or not text.strip():
        return False
    cjk = sum(1 for ch in text if '一' <= ch <= '鿿')
    non_space = sum(1 for ch in text if not ch.isspace())
    if non_space == 0:
        return False
    return (cjk / non_space) >= cjk_ratio_cutoff


async def _expand_query_to_en(query: str) -> str:
    """中文 query → 英文术语。失败/超时返回空串（调用方退回到原 query）。"""
    import asyncio
    cached = _query_expand_cache.get(query)
    if cached is not None:
        return cached

    try:
        from llm_client import client as _llm_client, MODEL as _LLM_MODEL
        resp = await asyncio.wait_for(
            _llm_client.chat.completions.create(
                model=_LLM_MODEL,
                messages=[
                    {"role": "system", "content": QUERY_EXPAND_PROMPT},
                    {"role": "user", "content": query},
                ],
                temperature=0.2,
                max_tokens=120,
            ),
            timeout=_QUERY_EXPAND_TIMEOUT,
        )
        expansion = (resp.choices[0].message.content or "").strip()
        # 去掉 LLM 偶尔加的 "答：" / 引号 / 末尾句号
        expansion = expansion.strip('"\'。.').replace("答：", "").strip()
        # 限长，避免极端输出污染 BM25
        if len(expansion) > 300:
            expansion = expansion[:300]
        _query_expand_cache[query] = expansion
        return expansion
    except Exception as e:
        print(f"[query expand] 失败：{e}（按原 query 检索）")
        _query_expand_cache[query] = ""  # 缓存失败结果避免反复重试
        return ""


async def search_async(query: str, top_k: int = TOP_K) -> List[Dict]:
    """异步检索：先把中文 query 扩展为 中文+英文术语，再走原 BM25 +「英文配额」+「相关性裁判」。

    - 若 query 不是以中文为主，或过长，则跳过扩展直接走原 search()
    - LLM 失败/超时不影响检索，自动退回原 query
    - top-k 内强制保留至少 min(2, top_k) 条英文 PDF 的全文 chunk（如果候选池里
      存在分数 >= EN_QUOTA_MIN_SCORE 的英文 chunk），确保模型能看到英文 PDF 的
      Methods/Results 段
    - 最后过一道 LLM 相关性裁判：BM25 分数高但内容跑题（例如问"市场现状"
      命中的全是化学成分论文）→ 返回 [] 触发拒答，避免模型硬拼
    """
    if not query or not query.strip():
        return []
    if len(query) > _QUERY_EXPAND_MAX_LEN or not _is_mostly_chinese(query):
        # 走原始 search 也要保证配额：拉大池子后强制分层
        pool_size = max(top_k * 10, 100)
        candidates = search(query, top_k=pool_size)
        results = _enforce_en_quota(candidates, top_k=top_k, min_en_full=min(2, top_k))
        results = _enforce_tier_quota(candidates, top_k=top_k)
    else:
        expansion = await _expand_query_to_en(query)
        if expansion:
            merged_query = f"{query} {expansion}"
            pool_size = max(top_k * 10, 100)
            candidates = search(merged_query, top_k=pool_size)
            quota = min(2, top_k)
            results = _enforce_en_quota(candidates, top_k=top_k, min_en_full=quota)
            # 在英文配额基础上再做分层配额；用同一 candidates 池
            results = _enforce_tier_quota(candidates, top_k=top_k)
        else:
            pool_size = max(top_k * 10, 100)
            candidates = search(query, top_k=pool_size)
            results = _enforce_tier_quota(candidates, top_k=top_k)

    if not results:
        return results

    # 相关性裁判：仅当 query 较短（关键词式）且不是闲聊式时检查，
    # 句子级长 query 一般歧义小、命中也更精准，省一次 LLM 调用
    if await _is_retrieval_irrelevant(query, results):
        return []

    # 官方权威文件优先：query 明显在问药典/团标时，强制 pin 相关权威文件
    results = _enforce_authoritative_pin(query, results, top_k=top_k)

    return results


# ============ 官方权威文件 pin（药典 / 团标优先）============
# 用户明确要求：只要提到"药典"，一律优先 2025 版药典条文；
# 提到"团体标准 / 保健食品"时优先团标 PDF。
# 光靠 BM25 排序不稳（xlsx 里堆积大量"药典"字样容易霸榜），
# 这里直接按文件名硬 pin：命中即置顶。
AUTHORITATIVE_FILES = {
    # (query 触发关键词元组, 目标文件名前缀)
    "pharmacopoeia_2025": (
        ("药典", "中国药典", "2025 版", "2025版", "2025 年版", "2025年版", "pharmacopoeia"),
        ("2025版中国药典—桑白皮.docx", "2025版中国药典—2351 真菌毒素测定法.docx"),
    ),
    "health_food_standard": (
        ("团体标准", "团标", "保健食品", "T/CNHFA"),
        ("保健食品用原料桑白皮团体标准.pdf",),
    ),
    # 桑属分子鉴定：郑铭志的序列清洗结论覆盖/推翻了原论文
    # （原论文声称 ITS2 可鉴别桑属近缘种，清洗后发现分不开鸡桑/华桑等）
    # 一旦 query 涉及分子鉴定，必须让清洗报告独占前排，否则模型会被
    # 熊永兴/沈烈行/刘宸浩等原论文的旧结论带偏。
    "mulberry_molecular_id": (
        ("ITS", "ITS2", "DNA条形码", "DNA 条形码", "条形码", "分子鉴定", "分子鉴别",
         "序列", "psbA", "matK", "rbcL",
         "鸡桑", "华桑", "蒙桑", "构树", "柘树", "桑属", "近缘种"),
        ("国内桑属物种DNA序列数据清洗与分析报告 (1).pptx",
         "基于《桑白皮及其混伪品的DNA条形码鉴定研究》对国内桑属物种ITS2序列数据清洗.docx",
         "国内桑属植物学序列比对结果.docx",
         "国内桑属中药鉴定学序列比对结果.docx"),
    ),
}


def _query_hits_authoritative(query: str) -> List[str]:
    """返回本次 query 命中的权威文件名列表（按定义顺序）。"""
    q_lower = query.lower()
    hits: List[str] = []
    for _, (keywords, files) in AUTHORITATIVE_FILES.items():
        if any(kw.lower() in q_lower for kw in keywords):
            for f in files:
                if f not in hits:
                    hits.append(f)
    return hits


def _enforce_authoritative_pin(query: str, results: List[Dict], top_k: int) -> List[Dict]:
    """query 命中权威触发词时，把对应权威文件的 chunk pin 到 top-k 前排。

    - 从当前 results 里挑最多 2 条属于权威文件的 chunk 移到最前（保留原分数排序）
    - 如果当前 results 里根本没有权威文件的 chunk，重新拉一次原始 search（大池），
      把权威文件的最高分 chunk 顶进来
    - 只做"前 2 位保底"，避免所有 top-k 都被权威文件霸占、把其他相关论文全挤掉
    """
    if not results:
        return results
    target_files = _query_hits_authoritative(query)
    if not target_files:
        return results

    target_set = set(target_files)

    # 从大池里捞所有权威 chunk（不仅仅是当前 results 里的）
    # —— 权威文件必须"独占前排"，不能被 BM25 数量更多的原论文稀释
    pool = search(query, top_k=200)
    pool_pinned = [r for r in pool if r.get("filename") in target_set]
    others_pool = [r for r in pool if r.get("filename") not in target_set]

    if not pool_pinned:
        return results

    pool_pinned.sort(key=lambda r: r.get("score", 0), reverse=True)
    # 前排配额 = min(命中权威 chunk 数, ceil(top_k/2), 5)
    # 至少让权威文件占据 top-k 的前一半，压过原论文的多数派
    lead_n = min(len(pool_pinned), max(1, top_k // 2 + 1), 5, top_k)
    lead = pool_pinned[:lead_n]
    lead_ids = {(r.get("filename"), r.get("chunk_index")) for r in lead}
    tail = [r for r in (results + others_pool)
            if (r.get("filename"), r.get("chunk_index")) not in lead_ids]
    return (lead + tail)[:top_k]


# ============ 相关性裁判（防止 BM25 高分但内容跑题）============
RELEVANCE_JUDGE_PROMPT = """你是中医药文献检索质量审核员。判断检索到的资料是否**有可能**回答用户问题。

判定原则（宽松判）：
- 只要资料中能找到与用户问题主题**直接相关**的信息（哪怕只是片段或表格数据），就 → YES
- 资料只是「沾了同样的关键词」但通篇讨论**完全不同的主题**才 → NO
  例如：
  · 用户问"市场现状/价格/产销"，资料全是化学成分/HPLC 检测 → NO
  · 用户问"某药材的真菌毒素污染"，资料是通用中药材真菌毒素检测方法综述、未涉及该药材 → NO
  · 用户问"产地加工/采收"，资料全是本草考证、产地分布 → 可考虑 NO
  · 用户问"性状/产地差异"，资料是不同产地的成分含量数据或品质评价 → 算 YES（成分差异属于性状差异的延伸）
  · 用户问"栽培 vs 野生"，资料里出现了两类来源的对比数据 → 算 YES

宁可放过，不要错杀。只输出一个英文单词：YES 或 NO。"""

_RELEVANCE_TIMEOUT = 3.0
_RELEVANCE_MAX_QUERY_LEN = 40   # query 超过此长度（句子级）跳过裁判
_RELEVANCE_JUDGE_TOP_N = 5      # 喂给裁判的 chunk 数（比 top-3 多看一些以减少误杀）
_RELEVANCE_SNIPPET_CHARS = 300  # 每条 chunk 的节选长度
_RELEVANCE_CACHE: Dict[str, bool] = {}


async def _is_retrieval_irrelevant(query: str, results: List[Dict]) -> bool:
    """True = 资料跑题应拒答；False = 资料对题继续走 RAG。
    失败/超时 fail-open，按"对题"处理，绝不阻塞正常回答。
    """
    import asyncio
    # 长 query 跳过：句子级问题 BM25 命中误差小，省一次 LLM 调用
    if len(query) > _RELEVANCE_MAX_QUERY_LEN:
        return False

    # 缓存 key：query + top-N filename，避免同一题反复判
    cache_key = query + "||" + "|".join(
        r.get("filename", "")[:50] for r in results[:_RELEVANCE_JUDGE_TOP_N]
    )
    if cache_key in _RELEVANCE_CACHE:
        return _RELEVANCE_CACHE[cache_key]

    # 构造裁判输入：title + 前 N 字
    snippets = []
    for i, r in enumerate(results[:_RELEVANCE_JUDGE_TOP_N], 1):
        title = r.get("title") or r.get("filename", "")[:60]
        text = r.get("text", "")[:_RELEVANCE_SNIPPET_CHARS].replace("\n", " ")
        snippets.append(f"【{i}】标题：{title}\n节选：{text}")
    user_msg = f"问题：{query}\n\n检索到的资料：\n" + "\n\n".join(snippets)

    try:
        from llm_client import client as _llm_client, MODEL as _LLM_MODEL
        resp = await asyncio.wait_for(
            _llm_client.chat.completions.create(
                model=_LLM_MODEL,
                messages=[
                    {"role": "system", "content": RELEVANCE_JUDGE_PROMPT},
                    {"role": "user", "content": user_msg},
                ],
                temperature=0.0,
                max_tokens=4,
            ),
            timeout=_RELEVANCE_TIMEOUT,
        )
        verdict = (resp.choices[0].message.content or "").strip().upper()
        irrelevant = verdict.startswith("NO")
        _RELEVANCE_CACHE[cache_key] = irrelevant
        if irrelevant:
            print(f"[relevance judge] 跑题拒答：{query!r}")
        return irrelevant
    except Exception as e:
        print(f"[relevance judge] 失败 fail-open：{e}")
        return False  # 失败时不拦截


# 英文 chunk 想"挤进" top-k，至少要满足这个 BM25 分数。
# 经验值：top-k 里平均分通常在 50-150 之间，这里设 25 大致是 top-k 中位的 1/3，
# 避免硬塞 score < 25 的低质 chunk 反而稀释 context。
EN_QUOTA_MIN_SCORE = 25.0


def _is_en_full_chunk(chunk: Dict) -> bool:
    """判断 chunk 是不是「英文 PDF 的全文片段」。

    标准（基于文件名最稳定，不会被段首公式/Figure 误导）：
    1. 不是 LLM 生成的中文摘要 chunk（开头不是【中文摘要】）
    2. 文件名里几乎不含中文（CJK 占比 < 20%）—— 这能精准锁定英文 PDF
    """
    text = chunk.get("text", "")
    if not text or text.startswith("【中文摘要】"):
        return False
    filename = chunk.get("filename", "")
    if not filename:
        return False
    cjk = sum(1 for ch in filename if '一' <= ch <= '鿿')
    total = len(filename)
    if total == 0:
        return False
    return (cjk / total) < 0.20


def _enforce_en_quota(candidates: List[Dict], top_k: int, min_en_full: int = 2) -> List[Dict]:
    """在按 BM25 分数排好序的候选里，强制保留至少 min_en_full 条英文 PDF 的全文 chunk。

    实现：先按分数填 top-k；若英文 full chunk 不足 min_en_full，从尾部把
    最弱的非英文 chunk 换成后续候选中分数最高的英文 chunk。
    只有英文 chunk 自身 score >= EN_QUOTA_MIN_SCORE 才会被插入，避免拉低质量。
    """
    if not candidates:
        return []
    selected = candidates[:top_k]
    en_in_selected = [c for c in selected if _is_en_full_chunk(c)]
    if len(en_in_selected) >= min_en_full:
        return selected

    need = min_en_full - len(en_in_selected)
    # 候选池里还没被选中的英文 chunk（按分数已排序），且通过质量阈值
    extra_en = [c for c in candidates[top_k:]
                if _is_en_full_chunk(c) and c["score"] >= EN_QUOTA_MIN_SCORE][:need]
    if not extra_en:
        return selected

    non_en_indices = [i for i, c in enumerate(selected) if not _is_en_full_chunk(c)]
    for en_chunk in extra_en:
        if not non_en_indices:
            break
        replace_idx = non_en_indices.pop()
        selected[replace_idx] = en_chunk
    selected.sort(key=lambda x: x["score"], reverse=True)
    return selected


# ============ 证据分层配额（A 类必召 + C 类封顶）============
# 目标：top-k 里至少含 min_a 条 A 类文献（若候选池里存在），C 类不超过 max_c
# —— 这是"方法学审查框架"的核心：让高证据等级文献稳定进入上下文
CRED_MIN_A_IN_TOPK = 1   # top-k 至少保留 1 条 A 类
CRED_MAX_C_IN_TOPK = 2   # top-k 最多容纳 2 条 C 类，避免低质文献占用配额


def _enforce_tier_quota(
    candidates: List[Dict],
    top_k: int,
    min_a: int = CRED_MIN_A_IN_TOPK,
    max_c: int = CRED_MAX_C_IN_TOPK,
) -> List[Dict]:
    """强制配额：至少 min_a 条 A 类进 top-k；C 类最多 max_c 条。

    - candidates 已按 final_score 从高到低排好
    - A 类不足时，从候选池后半段找最高分 A 类补上（前提：BM25 命中，说明确实相关）
    - C 类超额时，把超出的 C 类挤出 top-k，让位给候选池里下一位非 C 类
    - 若 query 本身命中的 A 类文献 BM25 为 0（完全跑题），不强塞（避免答非所问）
    """
    if not candidates or top_k <= 0:
        return candidates[:top_k]

    selected = candidates[:top_k]
    a_in_selected = [c for c in selected if c.get("tier") == "A"]
    c_in_selected = [c for c in selected if c.get("tier") == "C"]

    # ---- 1) A 类不足：从池子里补 ----
    if len(a_in_selected) < min_a:
        need = min_a - len(a_in_selected)
        # 候选池后半段里分数最高的 A 类（且要求 BM25 > 0，避免硬塞跑题内容）
        extra_a = [c for c in candidates[top_k:]
                   if c.get("tier") == "A" and c.get("bm25", 0) > 0][:need]
        if extra_a:
            # 从 selected 尾部替换最弱的非 A 项
            non_a_idx = [i for i, c in enumerate(selected) if c.get("tier") != "A"]
            for a_chunk in extra_a:
                if not non_a_idx:
                    break
                replace_i = non_a_idx.pop()  # pop 最后一个 = 分数最低
                selected[replace_i] = a_chunk

    # ---- 2) C 类超额：挤出多余 C 类 ----
    c_indices = [i for i, c in enumerate(selected) if c.get("tier") == "C"]
    if len(c_indices) > max_c:
        # 保留分数最高的 max_c 个 C 类，其余替换成候选池里下一个非 C 类
        c_sorted = sorted(c_indices, key=lambda i: selected[i]["score"], reverse=True)
        keep_c = set(c_sorted[:max_c])
        drop_c_indices = [i for i in c_indices if i not in keep_c]

        # 候选池后段的非 C 类补位
        used_ids = {(c.get("filename"), c.get("chunk_index")) for c in selected}
        backfill = [c for c in candidates[top_k:]
                    if c.get("tier") != "C" and (c.get("filename"), c.get("chunk_index")) not in used_ids]
        for drop_i in drop_c_indices:
            if not backfill:
                break
            selected[drop_i] = backfill.pop(0)

    # 重新按 score 排序输出
    selected.sort(key=lambda x: x["score"], reverse=True)
    return selected


# ============ BM25 检索 ============
def search(query: str, top_k: int = TOP_K) -> List[Dict]:
    """BM25 关键词检索，内存友好。"""
    index = load_index()
    chunks = index.get("chunks", [])
    doc_freq = index.get("doc_freq", {})
    total_docs = index.get("total_docs", 0)

    if not chunks or total_docs == 0:
        return []

    # BM25 参数
    k1 = 1.5
    b = 0.75

    # 计算平均文档长度
    avg_dl = sum(len(c["text"]) for c in chunks) / total_docs

    # 对 query 分词
    query_tokens = tokenize(query)
    if not query_tokens:
        return []

    # 计算每个文档的 BM25 分数，融合可信度加权
    scores = []
    for i, chunk in enumerate(chunks):
        doc_text = chunk["text"]
        # 打分 token 空间：正文 + filename + title（后两者软加权），
        # 与 build_index 阶段的 doc_freq 计算保持一致
        doc_tokens = _chunk_tokens(chunk)
        # doc_len 依旧用正文字符数，避免文件名极短时被 BM25 长度归一因子放大
        doc_len = len(doc_text)
        tf_counter = Counter(doc_tokens)

        bm25 = 0.0
        for token in query_tokens:
            if token not in tf_counter:
                continue
            tf = tf_counter[token]
            df = doc_freq.get(token, 0)
            if df == 0:
                continue
            # IDF
            idf = math.log((total_docs - df + 0.5) / (df + 0.5) + 1)
            # TF 归一化
            tf_norm = (tf * (k1 + 1)) / (tf + k1 * (1 - b + b * doc_len / avg_dl))
            bm25 += idf * tf_norm

        if bm25 > 0:
            # 可信度加权：final = bm25 * (BASE + (1-BASE) * credibility)
            cred = _cred_of(chunk.get("filename", ""))
            final = bm25 * (CRED_WEIGHT_BASE + (1 - CRED_WEIGHT_BASE) * cred)
            scores.append((final, bm25, cred, i))

    # 排序取 top_k（按加权后 final 排）
    scores.sort(reverse=True)
    results = []
    for final, bm25, cred, idx in scores[:top_k]:
        chunk = chunks[idx]
        results.append({
            "text": chunk["text"],
            "title": chunk.get("title", "未知"),
            "filename": chunk.get("filename", ""),
            "score": round(final, 4),        # 供后续排序/裁判用的实际得分
            "bm25": round(bm25, 4),          # 保留原始 BM25 供 debug
            "credibility": round(cred, 3),
            "tier": _tier_of_score(cred),
            "chunk_index": idx,
        })

    return results


def format_context_for_prompt(search_results: List[Dict], max_chars: int = 6000) -> str:
    """将检索结果格式化为可注入 system prompt 的参考文本。

    max_chars 默认 6000：deepseek 64K context 完全够用，3000 在 top_k=10 时
    截断太狠会让"深度思考多检索"白费。

    按 filename 聚合：同一篇论文的多个 chunk 合并为一条【参考X】，最多保留
    命中最高的 3 段，避免同一 filename 被列成"参考1/参考3/参考6"多个条目。
    """
    if not search_results:
        return ""

    # 按 filename 分组，保持首次出现顺序（=最高分先出现）
    grouped: Dict[str, Dict] = {}
    order: List[str] = []
    for item in search_results:
        fn = item.get("filename") or item.get("title") or "unknown"
        if fn not in grouped:
            # 优先从 search result 里带的 tier（已按当前阈值算好），兜底用 filename 反查
            tier = item.get("tier") or _tier_of_filename(fn)
            cred_rec = load_credibility().get(fn, {}) or {}
            grouped[fn] = {
                "title": item.get("title") or fn,
                "chunks": [],
                "tier": tier,
                "credibility": item.get("credibility", cred_rec.get("credibility", "?")),
                "doc_kind": cred_rec.get("doc_kind", "未标注"),
            }
            order.append(fn)
        if len(grouped[fn]["chunks"]) < 3:  # 每篇最多 3 段
            grouped[fn]["chunks"].append(item.get("text", ""))

    parts = []
    total_len = 0
    for i, fn in enumerate(order, 1):
        g = grouped[fn]
        merged_text = "\n---\n".join(g["chunks"])
        hit_note = f"（共命中 {len(g['chunks'])} 段）" if len(g["chunks"]) > 1 else ""
        cred_str = g["credibility"] if isinstance(g["credibility"], str) else f"{g['credibility']:.2f}"
        tier_tag = f"[{g['tier']}类·{g['doc_kind']}·可信度{cred_str}]"
        entry = f"【参考{i}】{tier_tag}（来源：{g['title']}）{hit_note}\n{merged_text}\n"
        if total_len + len(entry) > max_chars:
            break
        parts.append(entry)
        total_len += len(entry)

    if not parts:
        return ""

    header = (
        "以下是从知识库中收录的公开研究文献（涵盖历年桑白皮相关论文、《中国药典》条文、团体标准等）"
        "中检索到的相关内容。回答时必须以这些内容为依据，不得编造资料中没有的信息。\n\n"
        "【证据分层说明——本次回答的引用规则】\n"
        "- 每条【参考N】都标注了证据等级：**A 类**（药典/团标/SCI 顶级原始研究，可信度 ≥ 0.80）、"
        "**B 类**（核心期刊原始实验、方法学基本完整的综述，0.60–0.80）、**C 类**（普刊/科普/市场/老文献，< 0.60）。\n"
        "- **优先采纳 A 类证据**作为结论主干；B 类作为方法学细节与佐证；C 类**仅在 A/B 类都未涉及时**"
        "作为补充线索使用，且必须明确标注『该结论仅见于 C 类文献，尚待高等级证据验证』。\n"
        "- 若 A 类与 B/C 类结论冲突，**以 A 类为准**，同时以中性语言呈现 B/C 类的不同报道供参考。\n"
        "- **覆盖性反驳优先规则**：当参考资料里同时出现『某方法/序列可鉴别 X』与『经数据清洗/重新分析后该方法无法鉴别 X』时，"
        "一律以**后者（清洗/复核结论）**为准，前者视为已被推翻的旧结论。典型场景：桑属 ITS2 分子鉴定——"
        "早期原论文（如熊永兴、沈烈行、刘宸浩等的 DNA 条形码研究）声称 ITS2 可区分鸡桑/华桑/蒙桑等桑属近缘种，"
        "但郑铭志《国内桑属物种 DNA 序列数据清洗与分析报告》经序列清洗后证实：ITS2 只能分开构树、柘树等**其他属**，"
        "**分不开鸡桑、华桑等桑的近缘种**。回答此类问题必须以清洗结论为准，不得沿用原论文的旧说法。\n"
        "- 回答中每处关键数据/结论后，请用简短括注标明来源等级，如"
        "『（据《中国药典》2025 版，A 类）』『（相关 SCI 研究，B 类）』；无需完整点作者/年份的夹注。\n\n"
        "【定位说明——写作口吻硬规则】\n"
        "- 这些文献是**已公开发表的现有研究**，不是当前老师或团队产出的成果。\n"
        "- 讲述时一律以**第三人称、客观转述**的方式行文，例如：『现有研究表明……』"
        "『相关文献报道……』『已有研究显示……』『《中国药典》2025 年版一部规定……』。\n"
        "- **严禁**使用『我们的研究』『我团队』『本团队』『我们发现』『我们的实验』"
        "『我们的样品』这类第一人称集体表述——课堂讲解在转述行业已有工作，不是在汇报自己组的成果。\n"
        "- **严禁**在正文里塞入『X 团队 X 年报道』『某某某（2023）指出』『据〈作者，年份〉报道』"
        "『熊永兴团队 2013 年、刘宸浩团队 2026 年、沈烈行团队 2023 年的报道』等**夹注式引文标注**——"
        "课堂讲解不是论文写作，不用逐一点作者、点年份。直接把研究结论讲出来，"
        "只用『现有研究表明』『相关研究显示』『已有报道』这类中性引出词过渡即可。\n\n"
        "【关键纪律】检索系统按关键词重叠度排序，**沾关键词 ≠ 真正回答了用户的问题**。请你在动笔之前，\n"
        "先判断以下参考资料是否**直接**论述了用户所问的话题：\n"
        "  · 用户问『市场现状/价格/产销』而资料只讲化学成分、质量控制 → 属于跑题；\n"
        "  · 用户问『某药材的真菌毒素』而资料只是中药材通用真菌毒素综述、未具体研究该药材 → 属于跑题；\n"
        "  · 用户问『某具体加工步骤』而资料只讲产地分布、本草考证 → 属于跑题。\n"
        "遇到上述情况，请直接**按中药鉴定学的学科通识正面回答**用户问题——用规范术语讲原理、"
        "讲分类、讲鉴别要点，不许输出『资料较少』『知识库中没有』『暂时无法回答』『不便直接展开』这类"
        "搪塞话；参考资料里如果有部分沾边的信息，可以第三人称转述作为补充，但**严禁**把不切题的"
        "内容硬说成对用户问题的直接回答。核心原则：讲得到就正面讲，讲不到具体数据就讲学科通识，"
        "**永远不要用『资料不足』当遁词**。\n\n"
    )
    return header + "\n".join(parts)


# ============ 状态查询 ============
def get_index_stats() -> Dict:
    """获取索引状态统计"""
    try:
        index = load_index()
        chunks = index.get("chunks", [])
        if chunks:
            filenames = set(c.get("filename", "") for c in chunks)
            return {
                "total_chunks": len(chunks),
                "total_files": len(filenames),
                "files": sorted(filenames),
                "status": "ready",
            }
        else:
            return {
                "total_chunks": 0,
                "total_files": 0,
                "files": [],
                "status": "empty",
            }
    except Exception as e:
        return {
            "total_chunks": 0,
            "total_files": 0,
            "files": [],
            "status": f"error: {e}",
        }
